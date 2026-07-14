from __future__ import annotations

import base64
import hashlib
import io
import json
import os
import re
import shutil
import stat
import subprocess
import tempfile
import time
import uuid
import zipfile
from datetime import datetime, timezone
from dataclasses import dataclass
from pathlib import Path
from typing import Any

import magic
import requests
from defusedxml import ElementTree
from docx import Document
from openpyxl import load_workbook
from PIL import Image, ImageOps
from pillow_heif import register_heif_opener
from pptx import Presentation
from pypdf import PdfReader
from supabase import Client, create_client

register_heif_opener()
Image.MAX_IMAGE_PIXELS = 80_000_000

SUPABASE_URL = os.environ["SUPABASE_URL"]
SUPABASE_SERVICE_ROLE_KEY = os.environ["SUPABASE_SERVICE_ROLE_KEY"]
OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY", "")
STT_MODEL = os.getenv("ATTACHMENT_STT_MODEL", "openai/gpt-4o-mini-transcribe")
SUMMARY_MODEL = os.getenv("ATTACHMENT_OPENROUTER_MODEL", "google/gemini-2.5-flash")
POLL_SECONDS = max(0.5, float(os.getenv("ATTACHMENT_WORKER_POLL_SECONDS", "2")))
MAX_ATTEMPTS = max(1, int(os.getenv("ATTACHMENT_WORKER_MAX_ATTEMPTS", "4")))
WORKER_ID = os.getenv("RENDER_INSTANCE_ID") or f"attachment-worker-{uuid.uuid4()}"
ENVIRONMENT = os.getenv("ENVIRONMENT", "production").lower()
DEV_BYPASS = os.getenv("ATTACHMENT_SCANNER_DEV_BYPASS", "false").lower() == "true"
MAX_EXTRACTED_CHARS = 120_000
MAX_ARCHIVE_FILES = 500
MAX_ARCHIVE_EXPANDED_BYTES = 200 * 1024 * 1024
MAX_ARCHIVE_COMPRESSION_RATIO = 100
MAX_ARCHIVE_MEMBER_BYTES = 16 * 1024 * 1024
MAX_TEXT_FILE_BYTES = 16 * 1024 * 1024
TIKA_APP_PATH = os.getenv("TIKA_APP_PATH", "/opt/tika/tika-app.jar")

supabase: Client = create_client(SUPABASE_URL, SUPABASE_SERVICE_ROLE_KEY)

BLOCKED_EXTENSIONS = {
    "exe", "dll", "dylib", "so", "app", "dmg", "pkg", "msi", "apk", "jar", "com", "bat", "cmd", "ps1", "vbs", "scr",
    "docm", "dotm", "xlsm", "xltm", "xlam", "pptm", "potm", "ppam", "ppsm", "sldm",
}
NESTED_ARCHIVES = {"zip", "7z", "rar", "tar", "gz", "bz2", "xz"}


class RejectFile(Exception):
    pass


@dataclass
class Processed:
    text: str | None
    derivatives: list[dict[str, Any]]
    page_count: int | None = None
    duration_seconds: float | None = None
    width: int | None = None
    height: int | None = None


def log(event: str, **fields: Any) -> None:
    print(json.dumps({"event": event, "worker": WORKER_ID, **fields}, default=str), flush=True)


def now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


def table(name: str):
    return supabase.table(name)


def update_attachment(attachment_id: str, **patch: Any) -> None:
    table("chat_attachments").update(patch).eq("id", attachment_id).execute()


def update_job(job_id: str, **patch: Any) -> None:
    table("attachment_processing_jobs").update(patch).eq("id", job_id).execute()


def lease_job() -> dict[str, Any] | None:
    result = supabase.rpc("lease_attachment_processing_job", {
        "p_worker_id": WORKER_ID,
        "p_lease_seconds": 600,
    }).execute()
    rows = result.data or []
    return rows[0] if rows else None


def extension(name: str) -> str:
    return Path(name).suffix.lower().lstrip(".")


def run(command: list[str], timeout: int = 180) -> subprocess.CompletedProcess[str]:
    return subprocess.run(command, check=True, capture_output=True, text=True, timeout=timeout)


def scan(path: Path) -> str:
    if DEV_BYPASS and ENVIRONMENT in {"development", "local", "test"}:
        return "development bypass"
    if DEV_BYPASS:
        raise RuntimeError("ATTACHMENT_SCANNER_DEV_BYPASS is forbidden outside development")
    try:
        result = subprocess.run(
            ["clamdscan", "--fdpass", "--no-summary", str(path)],
            capture_output=True,
            text=True,
            timeout=120,
        )
    except (FileNotFoundError, subprocess.TimeoutExpired) as error:
        raise RuntimeError("Malware scanner is unavailable") from error
    detail = (result.stdout or result.stderr).strip()[:500]
    if result.returncode == 1:
        raise RejectFile("Malware scan rejected this file")
    if result.returncode != 0:
        raise RuntimeError(f"Malware scanner is unavailable: {detail or result.returncode}")
    return detail or "clean"


def verify_size(path: Path, expected: int) -> None:
    actual = path.stat().st_size
    if actual != expected:
        raise RejectFile(f"Uploaded size mismatch ({actual} != {expected})")
    if actual <= 0 or actual > 100 * 1024 * 1024:
        raise RejectFile("File size is outside the supported range")


def verify_magic(path: Path, row: dict[str, Any]) -> str:
    head = path.read_bytes()[:4096]
    if head.startswith(b"MZ") or head.startswith(b"\x7fELF") or head.startswith(b"\xca\xfe\xba\xbe"):
        raise RejectFile("Executable content is not supported")
    detected = magic.from_file(str(path), mime=True) or "application/octet-stream"
    ext = extension(row["original_name"])
    if ext in BLOCKED_EXTENSIONS:
        raise RejectFile("Executable and macro-enabled files are not supported")
    if zipfile.is_zipfile(path):
        detected = {
            "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "zip": "application/zip",
        }.get(ext, detected)
    declared = str(row.get("declared_mime_type") or "")
    declared_family = declared.split("/", 1)[0]
    detected_family = detected.split("/", 1)[0]
    if declared_family in {"image", "audio", "video"} and declared_family != detected_family:
        raise RejectFile("The declared MIME type does not match the file contents")
    return detected


def validate_zip(path: Path, office: bool = False) -> list[zipfile.ZipInfo]:
    with zipfile.ZipFile(path) as archive:
        infos = archive.infolist()
        if len(infos) > MAX_ARCHIVE_FILES:
            raise RejectFile("Archive contains too many files")
        total = 0
        names = set()
        for info in infos:
            normalized = info.filename.replace("\\", "/")
            pieces = normalized.split("/")
            if normalized.startswith("/") or ".." in pieces:
                raise RejectFile("Archive contains an unsafe path")
            mode = info.external_attr >> 16
            if stat.S_ISLNK(mode):
                raise RejectFile("Archive symlinks are not supported")
            if info.flag_bits & 0x1:
                raise RejectFile("Encrypted archives are not supported")
            ext = extension(normalized)
            if ext in BLOCKED_EXTENSIONS:
                raise RejectFile("Archive contains executable or macro-enabled content")
            if not office and ext in NESTED_ARCHIVES:
                raise RejectFile("Nested archives are not supported")
            total += info.file_size
            if total > MAX_ARCHIVE_EXPANDED_BYTES:
                raise RejectFile("Expanded archive is too large")
            if info.file_size > 1024 * 1024 and info.file_size / max(1, info.compress_size) > MAX_ARCHIVE_COMPRESSION_RATIO:
                raise RejectFile("Archive compression ratio is unsafe")
            names.add(normalized.lower())
        if any(name.endswith("vbaproject.bin") for name in names):
            raise RejectFile("Macro-enabled Office documents are not supported")
        if "encryptioninfo" in names or "encryptedpackage" in names:
            raise RejectFile("Password-protected Office documents are not supported")
        return infos


def validate_svg(path: Path) -> None:
    text = path.read_text("utf-8", errors="replace")
    if re.search(r"<script\b|<foreignObject\b|<!DOCTYPE|<!ENTITY|\son\w+\s*=|javascript:|data:text/html|@import|url\(\s*[\"']?\s*(?:https?:|file:|//)|(?:href|xlink:href)\s*=\s*[\"']\s*(?:https?:|file:|//)", text, re.I):
        raise RejectFile("SVG contains active or external content")
    try:
        root = ElementTree.fromstring(text)
    except Exception as error:
        raise RejectFile("Invalid SVG file") from error
    if not root.tag.lower().endswith("svg"):
        raise RejectFile("Invalid SVG file")


def read_text(path: Path) -> str:
    if path.stat().st_size > MAX_TEXT_FILE_BYTES:
        raise RejectFile("Text file is too large to process safely")
    return path.read_text("utf-8", errors="replace").replace("\x00", "").strip()[:MAX_EXTRACTED_CHARS]


def extract_pdf(path: Path) -> Processed:
    reader = PdfReader(str(path))
    if reader.is_encrypted:
        raise RejectFile("Password-protected PDFs are not supported")
    sections = []
    for index, page in enumerate(reader.pages[:2000], start=1):
        sections.append(f"[Page {index}]\n{(page.extract_text() or '').strip()}")
    return Processed("\n\n".join(sections).strip()[:MAX_EXTRACTED_CHARS], [], page_count=len(reader.pages))


def extract_docx(path: Path) -> Processed:
    validate_zip(path, office=True)
    document = Document(str(path))
    text = "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text.strip())
    for table_index, document_table in enumerate(document.tables, start=1):
        text += f"\n\n[Table {table_index}]\n"
        text += "\n".join("\t".join(cell.text for cell in row.cells) for row in document_table.rows)
    return Processed(text.strip()[:MAX_EXTRACTED_CHARS], [])


def extract_pptx(path: Path) -> Processed:
    validate_zip(path, office=True)
    deck = Presentation(str(path))
    sections = []
    for index, slide in enumerate(deck.slides, start=1):
        strings = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                strings.append(shape.text.strip())
        sections.append(f"[Slide {index}]\n" + "\n".join(strings))
    return Processed("\n\n".join(sections).strip()[:MAX_EXTRACTED_CHARS], [], page_count=len(deck.slides))


def extract_xlsx(path: Path) -> Processed:
    validate_zip(path, office=True)
    workbook = load_workbook(str(path), read_only=True, data_only=True)
    sections = []
    for sheet in workbook.worksheets[:100]:
        rows = []
        for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            if row_index > 5000:
                rows.append("[Rows truncated after 5000]")
                break
            rows.append("\t".join("" if value is None else str(value) for value in row[:100]))
        sections.append(f"[Sheet {sheet.title}]\n" + "\n".join(rows))
    return Processed("\n\n".join(sections).strip()[:MAX_EXTRACTED_CHARS], [])


def extract_legacy_office(path: Path, workdir: Path) -> Processed:
    outdir = workdir / f"converted-{uuid.uuid4().hex}"
    outdir.mkdir()
    run(["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", str(outdir), str(path)], timeout=240)
    converted = next(outdir.glob("*.pdf"), None)
    if not converted:
        raise RuntimeError("LibreOffice could not convert the document")
    return extract_pdf(converted)


def extract_with_tika(path: Path) -> Processed:
    if not Path(TIKA_APP_PATH).is_file():
        raise RuntimeError("Apache Tika is unavailable")
    result = run(["java", "-Xmx512m", "-jar", TIKA_APP_PATH, "-t", str(path)], timeout=240)
    text = result.stdout.replace("\x00", "").strip()
    if not text:
        raise RuntimeError("Apache Tika returned no extractable content")
    return Processed(text[:MAX_EXTRACTED_CHARS], [{"kind": "extraction", "engine": "apache-tika"}])


def upload_derivative(row: dict[str, Any], payload: bytes, filename: str, mime_type: str, kind: str, **metadata: Any) -> dict[str, Any]:
    base = row["storage_path"].rsplit("/", 1)[0]
    path = f"{base}/derivatives/{row['id']}-{filename}"
    supabase.storage.from_(row["bucket"]).upload(
        path,
        payload,
        file_options={"content-type": mime_type, "upsert": "true", "cache-control": "31536000"},
    )
    return {"kind": kind, "storage_path": path, "mimeType": mime_type, **metadata}


def describe_image(payload: bytes) -> str:
    if not OPENROUTER_API_KEY:
        if ENVIRONMENT in {"development", "local", "test"}:
            return ""
        raise RuntimeError("OPENROUTER_API_KEY is required for image fallback descriptions")
    response = requests.post(
        "https://openrouter.ai/api/v1/chat/completions",
        headers={"Authorization": f"Bearer {OPENROUTER_API_KEY}", "Content-Type": "application/json", "HTTP-Referer": "https://polyphonic.chat", "X-Title": "Polyphonic Attachments"},
        json={
            "model": SUMMARY_MODEL,
            "messages": [{
                "role": "user",
                "content": [
                    {"type": "text", "text": "Describe this uploaded image faithfully for an agent that cannot see it. Transcribe all legible text, preserve layout relationships that matter, and state uncertainty. Do not infer facts that are not visible."},
                    {"type": "image_url", "image_url": {"url": "data:image/webp;base64," + base64.b64encode(payload).decode()}},
                ],
            }],
            "temperature": 0.1,
            "max_tokens": 1600,
        },
        timeout=180,
    )
    response.raise_for_status()
    choice = (response.json().get("choices") or [{}])[0]
    if choice.get("finish_reason") != "stop":
        raise RuntimeError("Image description did not complete")
    description = str(choice.get("message", {}).get("content") or "").strip()
    if not description:
        raise RuntimeError("Image description was empty")
    return description


def process_image(path: Path, row: dict[str, Any], workdir: Path) -> Processed:
    ext = extension(row["original_name"])
    if ext == "svg":
        validate_svg(path)
        import cairosvg
        png_path = workdir / "safe.png"
        cairosvg.svg2png(url=str(path), write_to=str(png_path), output_width=2400, output_height=2400)
        source = Image.open(png_path)
    else:
        source = Image.open(path)
    source.seek(0)
    source.load()
    source = ImageOps.exif_transpose(source).convert("RGB")
    width, height = source.size
    safe = source.copy()
    safe.thumbnail((2400, 2400), Image.Resampling.LANCZOS)
    thumb = source.copy()
    thumb.thumbnail((640, 640), Image.Resampling.LANCZOS)
    safe_bytes = io.BytesIO()
    safe.save(safe_bytes, format="WEBP", quality=88, method=6)
    thumb_bytes = io.BytesIO()
    thumb.save(thumb_bytes, format="WEBP", quality=82, method=6)
    display_payload = safe_bytes.getvalue()
    derivatives = [
        upload_derivative(row, display_payload, "display.webp", "image/webp", "safe-display", width=safe.width, height=safe.height),
        upload_derivative(row, thumb_bytes.getvalue(), "thumb.webp", "image/webp", "thumbnail", width=thumb.width, height=thumb.height),
    ]
    description = describe_image(display_payload)
    if description:
        derivatives.append({"kind": "summary", "text": description, "engine": "openrouter-vision"})
    extracted = f"[Image description]\n{description}" if description else None
    return Processed(extracted, derivatives, width=width, height=height)


def ffprobe_duration(path: Path) -> float:
    result = run(["ffprobe", "-v", "error", "-show_entries", "format=duration", "-of", "default=noprint_wrappers=1:nokey=1", str(path)])
    return max(0.0, float(result.stdout.strip() or 0))


def timestamp(seconds: float) -> str:
    total = max(0, int(seconds))
    hours, remainder = divmod(total, 3600)
    minutes, secs = divmod(remainder, 60)
    return f"{hours:02d}:{minutes:02d}:{secs:02d}" if hours else f"{minutes:02d}:{secs:02d}"


def transcribe_audio(path: Path, workdir: Path) -> tuple[str, list[dict[str, Any]]]:
    if not OPENROUTER_API_KEY:
        raise RuntimeError("OPENROUTER_API_KEY is required for audio and video processing")
    segment_pattern = workdir / "audio-%03d.mp3"
    run([
        "ffmpeg", "-nostdin", "-y", "-i", str(path), "-vn", "-ac", "1", "-ar", "16000", "-b:a", "64k",
        "-f", "segment", "-segment_time", "600", "-reset_timestamps", "1", str(segment_pattern),
    ], timeout=600)
    segments = sorted(workdir.glob("audio-*.mp3"))
    if not segments:
        raise RuntimeError("No audio track could be extracted")
    transcript_parts = []
    timeline = []
    cursor = 0.0
    for segment in segments:
        duration = ffprobe_duration(segment)
        response = requests.post(
            "https://openrouter.ai/api/v1/audio/transcriptions",
            headers={"Authorization": f"Bearer {OPENROUTER_API_KEY}", "Content-Type": "application/json", "HTTP-Referer": "https://polyphonic.chat", "X-Title": "Polyphonic Attachments"},
            json={"model": STT_MODEL, "input_audio": {"data": base64.b64encode(segment.read_bytes()).decode(), "format": "mp3"}, "temperature": 0},
            timeout=120,
        )
        response.raise_for_status()
        text = str(response.json().get("text") or "").strip()
        start, end = cursor, cursor + duration
        if text:
            transcript_parts.append(f"[{timestamp(start)}–{timestamp(end)}] {text}")
            timeline.append({"start": start, "end": end, "text": text})
        cursor = end
    return "\n\n".join(transcript_parts), timeline


def summarize_video(transcript: str, keyframes: list[Path]) -> str:
    if not OPENROUTER_API_KEY:
        raise RuntimeError("OPENROUTER_API_KEY is required for video processing")
    content: list[dict[str, Any]] = [{
        "type": "text",
        "text": "Create a grounded scene summary of this uploaded video. Use only visible frames and the timestamped transcript. Mention uncertainty and organize by time.\n\nTranscript:\n" + transcript[:60_000],
    }]
    for frame in keyframes[:8]:
        content.append({"type": "image_url", "image_url": {"url": "data:image/jpeg;base64," + base64.b64encode(frame.read_bytes()).decode()}})
    response = requests.post(
        "https://openrouter.ai/api/v1/chat/completions",
        headers={"Authorization": f"Bearer {OPENROUTER_API_KEY}", "Content-Type": "application/json", "HTTP-Referer": "https://polyphonic.chat", "X-Title": "Polyphonic Attachments"},
        json={"model": SUMMARY_MODEL, "messages": [{"role": "user", "content": content}], "temperature": 0.2, "max_tokens": 1600},
        timeout=180,
    )
    response.raise_for_status()
    choice = (response.json().get("choices") or [{}])[0]
    if choice.get("finish_reason") != "stop":
        raise RuntimeError("Video scene summary did not complete")
    summary = str(choice.get("message", {}).get("content") or "").strip()
    if not summary:
        raise RuntimeError("Video scene summary was empty")
    return summary


def process_media(path: Path, row: dict[str, Any], workdir: Path) -> Processed:
    duration = ffprobe_duration(path)
    transcript, timeline = transcribe_audio(path, workdir)
    derivatives: list[dict[str, Any]] = [{"kind": "transcript", "text": transcript, "segments": timeline}]
    if row["kind"] == "audio":
        return Processed(transcript[:MAX_EXTRACTED_CHARS], derivatives, duration_seconds=duration)

    frame_dir = workdir / "frames"
    frame_dir.mkdir()
    interval = max(1.0, duration / 7.0) if duration else 10.0
    run([
        "ffmpeg", "-nostdin", "-y", "-i", str(path), "-vf", f"fps=1/{interval},scale='min(1280,iw)':-2", "-frames:v", "8", str(frame_dir / "frame-%02d.jpg")
    ], timeout=600)
    frames = sorted(frame_dir.glob("frame-*.jpg"))
    for index, frame in enumerate(frames):
        at = min(duration, interval * (index + 1))
        derivatives.append(upload_derivative(row, frame.read_bytes(), f"frame-{index + 1:02d}.jpg", "image/jpeg", "keyframe", timestampStart=at, label=timestamp(at)))
    summary = summarize_video(transcript, frames)
    if summary:
        derivatives.append({"kind": "summary", "text": summary})
    text = f"[Transcript]\n{transcript}\n\n[Scene summary]\n{summary}".strip()
    return Processed(text[:MAX_EXTRACTED_CHARS], derivatives, duration_seconds=duration)


def process_archive(path: Path, workdir: Path) -> Processed:
    infos = validate_zip(path, office=False)
    sections: list[str] = []
    text_extensions = {
        "txt", "md", "markdown", "csv", "tsv", "json", "jsonl", "xml", "html", "htm",
        "css", "scss", "js", "jsx", "ts", "tsx", "py", "rb", "go", "rs", "java", "c", "h",
        "cpp", "hpp", "cs", "sh", "bash", "zsh", "sql", "yaml", "yml", "toml", "ini",
    }
    extracted_chars = 0
    with zipfile.ZipFile(path) as archive:
        for index, info in enumerate(infos):
            if info.is_dir():
                continue
            header = f"[Archive member: {info.filename}]"
            if info.file_size > MAX_ARCHIVE_MEMBER_BYTES:
                section = f"{header}\n[Skipped: member exceeds {MAX_ARCHIVE_MEMBER_BYTES // (1024 * 1024)} MB extraction limit]"
                sections.append(section)
                extracted_chars += len(section)
                continue
            ext = extension(info.filename)
            try:
                payload = archive.read(info)
                if ext in text_extensions:
                    member_text = payload.decode("utf-8", errors="replace").replace("\x00", "").strip()
                    section = f"{header}\n{member_text}"
                    sections.append(section)
                elif ext in {"pdf", "docx", "pptx", "xlsx", "doc", "ppt", "xls", "rtf"}:
                    member_path = workdir / f"archive-member-{index}.{ext}"
                    member_path.write_bytes(payload)
                    if ext == "pdf":
                        processed = extract_pdf(member_path)
                    elif ext == "docx":
                        processed = extract_docx(member_path)
                    elif ext == "pptx":
                        processed = extract_pptx(member_path)
                    elif ext == "xlsx":
                        processed = extract_xlsx(member_path)
                    elif ext in {"doc", "ppt", "xls"}:
                        processed = extract_legacy_office(member_path, workdir)
                    else:
                        processed = extract_with_tika(member_path)
                    section = f"{header}\n{processed.text or '[No extractable text]'}"
                    sections.append(section)
                else:
                    section = f"{header}\n[Binary member: {info.file_size} bytes]"
                    sections.append(section)
                extracted_chars += len(section)
            except RejectFile:
                raise
            except Exception as error:
                section = f"{header}\n[Extraction unavailable: {str(error)[:240]}]"
                sections.append(section)
                extracted_chars += len(section)
            if extracted_chars >= MAX_EXTRACTED_CHARS:
                sections.append("[Archive extraction truncated at model context limit]")
                break
    return Processed("\n\n".join(sections)[:MAX_EXTRACTED_CHARS], [])


def process_file(path: Path, row: dict[str, Any], workdir: Path) -> Processed:
    kind = row["kind"]
    ext = extension(row["original_name"])
    if kind == "image":
        return process_image(path, row, workdir)
    if kind in {"text", "code"} or ext in {"csv", "tsv", "xml", "json"}:
        return Processed(read_text(path), [])
    if ext == "pdf":
        return extract_pdf(path)
    if ext == "docx":
        return extract_docx(path)
    if ext == "pptx":
        return extract_pptx(path)
    if ext == "xlsx":
        return extract_xlsx(path)
    if ext in {"doc", "ppt", "xls"}:
        return extract_legacy_office(path, workdir)
    if ext in {"rtf", "html", "htm"} or kind in {"document", "spreadsheet", "presentation"}:
        return extract_with_tika(path)
    if kind in {"audio", "video"}:
        return process_media(path, row, workdir)
    if kind == "archive":
        return process_archive(path, workdir)
    return Processed(None, [])


def process_job(job: dict[str, Any]) -> None:
    attachment_id = job["attachment_id"]
    rows = table("chat_attachments").select("*").eq("id", attachment_id).limit(1).execute().data or []
    if not rows:
        update_job(job["id"], status="cancelled", last_error="Attachment was deleted", locked_at=None, locked_by=None)
        return
    row = rows[0]
    if row["status"] == "cancelled":
        update_job(job["id"], status="cancelled", locked_at=None, locked_by=None)
        return

    with tempfile.TemporaryDirectory(prefix="polyphonic-attachment-") as temp:
        workdir = Path(temp)
        source = workdir / (Path(row["original_name"]).name or "attachment")
        payload = supabase.storage.from_(row["bucket"]).download(row["storage_path"])
        source.write_bytes(payload)
        verify_size(source, int(row["size_bytes"]))
        verified_mime = verify_magic(source, row)
        update_attachment(attachment_id, status="scanning", verified_mime_type=verified_mime, processing_error=None)
        scan_detail = scan(source)
        update_attachment(attachment_id, status="extracting", scanned_at=now_iso())
        processed = process_file(source, row, workdir)
        checksum = hashlib.sha256(source.read_bytes()).hexdigest()
        duplicates = table("chat_attachments").select("id").eq("user_id", row["user_id"]).eq("sha256", checksum).eq("status", "ready").neq("id", attachment_id).limit(1).execute().data or []
        derivatives = processed.derivatives + [{"kind": "scan", "engine": "clamav", "detail": scan_detail}]
        patch: dict[str, Any] = {
            "status": "ready",
            "verified_mime_type": verified_mime,
            "sha256": checksum,
            "duplicate_of": duplicates[0]["id"] if duplicates else None,
            "extracted_text": processed.text,
            "derivatives": derivatives,
            "page_count": processed.page_count,
            "duration_seconds": processed.duration_seconds,
            "width": processed.width,
            "height": processed.height,
            "ready_at": now_iso(),
            "processing_error": None,
        }
        update_attachment(attachment_id, **patch)
        update_job(job["id"], status="complete", last_error=None, locked_at=None, locked_by=None)
        log("attachment_ready", attachment_id=attachment_id, kind=row["kind"], size=row["size_bytes"], duplicate_of=patch["duplicate_of"])


def handle_failure(job: dict[str, Any], error: Exception) -> None:
    attachment_id = job["attachment_id"]
    message = str(error)[:1000]
    if isinstance(error, RejectFile):
        update_attachment(attachment_id, status="rejected", processing_error=message)
        rows = table("chat_attachments").select("bucket,storage_path").eq("id", attachment_id).limit(1).execute().data or []
        if rows:
            try:
                supabase.storage.from_(rows[0]["bucket"]).remove([rows[0]["storage_path"]])
            except Exception:
                pass
        update_job(job["id"], status="failed", last_error=message, locked_at=None, locked_by=None)
        log("attachment_rejected", attachment_id=attachment_id, error=message)
        return

    attempts = int(job.get("attempts") or 1)
    if attempts >= MAX_ATTEMPTS:
        update_attachment(attachment_id, status="failed", processing_error=message)
        update_job(job["id"], status="failed", last_error=message, locked_at=None, locked_by=None)
    else:
        delay = min(300, 10 * (2 ** max(0, attempts - 1)))
        available = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime(time.time() + delay))
        update_attachment(attachment_id, status="quarantined", processing_error=f"Processing retry {attempts}/{MAX_ATTEMPTS}: {message}")
        update_job(job["id"], status="queued", available_at=available, last_error=message, locked_at=None, locked_by=None)
    log("attachment_processing_failed", attachment_id=attachment_id, attempts=attempts, error=message)


def cleanup_orphaned_drafts() -> None:
    cutoff = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime(time.time() - 24 * 3600))
    try:
        table("chat_attachments").delete().is_("thread_id", "null").is_("room_id", "null").is_("message_id", "null").lt("created_at", cutoff).execute()
    except Exception as error:
        log("draft_cleanup_failed", error=str(error))


def main() -> None:
    log("worker_started", environment=ENVIRONMENT, scanner_bypass=DEV_BYPASS)
    last_cleanup = 0.0
    while True:
        try:
            if time.time() - last_cleanup > 3600:
                cleanup_orphaned_drafts()
                last_cleanup = time.time()
            job = lease_job()
            if not job:
                time.sleep(POLL_SECONDS)
                continue
            try:
                process_job(job)
            except Exception as error:
                handle_failure(job, error)
        except Exception as error:
            log("worker_loop_error", error=str(error))
            time.sleep(max(2.0, POLL_SECONDS))


if __name__ == "__main__":
    main()
